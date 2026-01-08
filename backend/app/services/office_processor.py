"""
Office Document Processor
Handles text extraction from Office documents (PowerPoint, Excel, Word)
"""

import logging
import io
import base64
from typing import Optional, Dict, Any, List
from datetime import datetime

logger = logging.getLogger(__name__)


class OfficeDocumentProcessor:
    """
    Processor for Microsoft Office documents (PowerPoint, Excel, Word).
    Extracts text content for RAG indexing and search.
    """
    
    def __init__(self):
        self._check_dependencies()
    
    def _check_dependencies(self):
        """Check if required libraries are available"""
        self.pptx_available = False
        self.xlsx_available = False
        self.docx_available = False
        
        try:
            from pptx import Presentation
            self.pptx_available = True
            logger.info("✅ python-pptx available for PowerPoint processing")
        except ImportError:
            logger.warning("⚠️ python-pptx not installed - PowerPoint extraction will be limited")
        
        try:
            from openpyxl import load_workbook
            self.xlsx_available = True
            logger.info("✅ openpyxl available for Excel processing")
        except ImportError:
            logger.warning("⚠️ openpyxl not installed - Excel extraction will be limited")
        
        try:
            from docx import Document
            self.docx_available = True
            logger.info("✅ python-docx available for Word processing")
        except ImportError:
            logger.warning("⚠️ python-docx not installed - Word extraction will be limited")
    
    def extract_text_from_bytes(self, file_bytes: bytes, file_type: str, filename: str = "") -> str:
        """
        Extract text from Office document bytes based on file type.
        
        Args:
            file_bytes: Raw document bytes
            file_type: MIME type of the document
            filename: Original filename (for extension-based detection)
            
        Returns:
            Extracted text content
        """
        file_type_lower = file_type.lower()
        filename_lower = filename.lower()
        
        # PowerPoint
        if 'presentation' in file_type_lower or 'powerpoint' in file_type_lower or \
           filename_lower.endswith('.pptx') or filename_lower.endswith('.ppt'):
            return self.extract_from_powerpoint(file_bytes)
        
        # Excel
        if 'spreadsheet' in file_type_lower or 'excel' in file_type_lower or \
           filename_lower.endswith('.xlsx') or filename_lower.endswith('.xls') or \
           filename_lower.endswith('.csv'):
            return self.extract_from_excel(file_bytes, filename)
        
        # Word
        if 'word' in file_type_lower or 'msword' in file_type_lower or \
           filename_lower.endswith('.docx') or filename_lower.endswith('.doc'):
            return self.extract_from_word(file_bytes)
        
        # Text/CSV
        if 'text' in file_type_lower or 'csv' in file_type_lower:
            return self.extract_from_text(file_bytes)
        
        logger.warning(f"Unknown Office file type: {file_type} ({filename})")
        return ""
    
    def extract_from_powerpoint(self, file_bytes: bytes) -> str:
        """
        Extract text from PowerPoint file.
        
        Args:
            file_bytes: PowerPoint file bytes
            
        Returns:
            Extracted text from all slides
        """
        if not self.pptx_available:
            logger.warning("python-pptx not available, cannot extract PowerPoint text")
            return ""
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation(io.BytesIO(file_bytes))
            text_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_texts = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                    
                    # Handle tables
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text:
                                    row_text.append(cell.text)
                            if row_text:
                                slide_texts.append(" | ".join(row_text))
                
                if slide_texts:
                    text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"📊 Extracted {len(extracted_text)} chars from PowerPoint ({len(prs.slides)} slides)")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error extracting PowerPoint text: {e}")
            return ""
    
    def extract_from_excel(self, file_bytes: bytes, filename: str = "") -> str:
        """
        Extract text from Excel file.
        
        Args:
            file_bytes: Excel file bytes
            filename: Original filename
            
        Returns:
            Extracted text from all sheets
        """
        # Handle CSV files
        if filename.lower().endswith('.csv'):
            return self.extract_from_csv(file_bytes)
        
        if not self.xlsx_available:
            logger.warning("openpyxl not available, cannot extract Excel text")
            return ""
        
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            text_parts = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_texts = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_values = []
                    for cell in row:
                        if cell is not None:
                            row_values.append(str(cell))
                    if row_values:
                        sheet_texts.append(" | ".join(row_values))
                
                if sheet_texts:
                    text_parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_texts))
            
            wb.close()
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"📊 Extracted {len(extracted_text)} chars from Excel ({len(wb.sheetnames)} sheets)")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error extracting Excel text: {e}")
            return ""
    
    def extract_from_csv(self, file_bytes: bytes) -> str:
        """
        Extract text from CSV file.
        
        Args:
            file_bytes: CSV file bytes
            
        Returns:
            Extracted text
        """
        try:
            import csv
            
            # Try different encodings
            text = None
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    text = file_bytes.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            
            if not text:
                logger.warning("Could not decode CSV with any encoding")
                return ""
            
            # Parse CSV
            reader = csv.reader(io.StringIO(text))
            rows = []
            for row in reader:
                if row:
                    rows.append(" | ".join(row))
            
            extracted_text = "\n".join(rows)
            logger.info(f"📊 Extracted {len(extracted_text)} chars from CSV ({len(rows)} rows)")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error extracting CSV text: {e}")
            return ""
    
    def extract_from_word(self, file_bytes: bytes) -> str:
        """
        Extract text from Word document.
        
        Args:
            file_bytes: Word document bytes
            
        Returns:
            Extracted text
        """
        if not self.docx_available:
            logger.warning("python-docx not available, cannot extract Word text")
            return ""
        
        try:
            from docx import Document
            
            doc = Document(io.BytesIO(file_bytes))
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                table_texts = []
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        table_texts.append(" | ".join(row_text))
                if table_texts:
                    text_parts.append("\n".join(table_texts))
            
            extracted_text = "\n\n".join(text_parts)
            logger.info(f"📄 Extracted {len(extracted_text)} chars from Word document")
            return extracted_text
            
        except Exception as e:
            logger.error(f"Error extracting Word text: {e}")
            return ""
    
    def extract_from_text(self, file_bytes: bytes) -> str:
        """
        Extract text from plain text files.
        
        Args:
            file_bytes: Text file bytes
            
        Returns:
            Decoded text
        """
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252', 'utf-16']:
                try:
                    text = file_bytes.decode(encoding)
                    logger.info(f"📝 Extracted {len(text)} chars from text file ({encoding})")
                    return text
                except UnicodeDecodeError:
                    continue
            
            logger.warning("Could not decode text file with any encoding")
            return ""
            
        except Exception as e:
            logger.error(f"Error extracting text: {e}")
            return ""
    
    def get_document_metadata(self, file_bytes: bytes, file_type: str, filename: str = "") -> Dict[str, Any]:
        """
        Get metadata from Office document.
        
        Args:
            file_bytes: Document bytes
            file_type: MIME type
            filename: Original filename
            
        Returns:
            Document metadata
        """
        metadata = {
            "filename": filename,
            "file_type": file_type,
            "size_bytes": len(file_bytes),
            "extracted_at": datetime.now().isoformat()
        }
        
        file_type_lower = file_type.lower()
        
        try:
            # PowerPoint metadata
            if 'presentation' in file_type_lower or 'powerpoint' in file_type_lower:
                if self.pptx_available:
                    from pptx import Presentation
                    prs = Presentation(io.BytesIO(file_bytes))
                    metadata["slide_count"] = len(prs.slides)
                    if prs.core_properties:
                        metadata["title"] = prs.core_properties.title
                        metadata["author"] = prs.core_properties.author
                        metadata["created"] = str(prs.core_properties.created) if prs.core_properties.created else None
            
            # Excel metadata
            elif 'spreadsheet' in file_type_lower or 'excel' in file_type_lower:
                if self.xlsx_available:
                    from openpyxl import load_workbook
                    wb = load_workbook(io.BytesIO(file_bytes), read_only=True)
                    metadata["sheet_count"] = len(wb.sheetnames)
                    metadata["sheet_names"] = wb.sheetnames
                    wb.close()
            
            # Word metadata
            elif 'word' in file_type_lower or 'msword' in file_type_lower:
                if self.docx_available:
                    from docx import Document
                    doc = Document(io.BytesIO(file_bytes))
                    metadata["paragraph_count"] = len(doc.paragraphs)
                    metadata["table_count"] = len(doc.tables)
                    if doc.core_properties:
                        metadata["title"] = doc.core_properties.title
                        metadata["author"] = doc.core_properties.author
        
        except Exception as e:
            logger.warning(f"Could not extract metadata: {e}")
        
        return metadata


# Global instance for easy access
office_processor = OfficeDocumentProcessor()
